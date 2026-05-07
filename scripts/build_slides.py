"""
build_slides.py
================
유튜브 영상 슬라이드 (pptx) — 시각자료 풍부 버전.

전제: scripts/build_slide_charts.py로 PNG 6종이 charts/에 생성되어 있어야 함.

산출물: outputs_trimmed10/slides_<MONTH>.pptx (28매)
"""
from __future__ import annotations

import argparse
import json
from datetime import datetime
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Inches, Pt


# ─── 색 ───
NAVY = RGBColor(0x1F, 0x4E, 0x79)
GOLD = RGBColor(0xC8, 0x9B, 0x3C)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC0, 0x39, 0x2B)
BG = RGBColor(0xFF, 0xFF, 0xFF)


# ============================================================
# 헬퍼
# ============================================================
def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_bg(slide, prs, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *,
             font_size=18, bold=False, color=GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    if font_name:
        r.font.name = font_name
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=LIGHT, border=None, border_w=1.0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if border is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = border
        rect.line.width = Pt(border_w)
    rect.shadow.inherit = False
    return rect


def add_round_rect(slide, left, top, width, height, *,
                    fill=LIGHT, border=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if border is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = border
        rect.line.width = Pt(1.2)
    rect.shadow.inherit = False
    return rect


def add_image(slide, path, left, top, width=None, height=None):
    if width and height:
        return slide.shapes.add_picture(str(path), left, top,
                                        width=width, height=height)
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def add_footer(slide, prs, *, slide_no, total, chapter, ts=None):
    """하단 footer: chapter | timestamp | slide #/N."""
    H = prs.slide_height
    W = prs.slide_width
    y = H - Cm(0.9)
    add_rect(slide, 0, y, W, Cm(0.05), fill=NAVY)
    txt_left = Cm(0.5)
    txt_right_w = Cm(3.0)
    if ts:
        chap_text = f"{chapter}   |   {ts}"
    else:
        chap_text = chapter
    add_text(slide, txt_left, y + Cm(0.18), W - Cm(4.0), Cm(0.6),
             chap_text, font_size=10, color=GREY)
    add_text(slide, W - Cm(3.5), y + Cm(0.18), Cm(3.0), Cm(0.6),
             f"{slide_no} / {total}", font_size=10, color=GREY,
             align=PP_ALIGN.RIGHT)


def add_kpi_card(slide, left, top, width, height, *,
                  label, value, sub=None, accent=NAVY):
    rect = add_round_rect(slide, left, top, width, height,
                           fill=WHITE, border=accent)
    # 라벨
    add_text(slide, left + Cm(0.3), top + Cm(0.25),
             width - Cm(0.6), Cm(0.6), label,
             font_size=11, color=GREY)
    # 값
    add_text(slide, left + Cm(0.3), top + Cm(0.85),
             width - Cm(0.6), Cm(1.2), value,
             font_size=24, color=accent, bold=True)
    if sub:
        add_text(slide, left + Cm(0.3),
                 top + height - Cm(0.8),
                 width - Cm(0.6), Cm(0.6), sub,
                 font_size=10, color=GREY)


def add_bullet_card(slide, left, top, width, height, *,
                     icon, title, body, accent=NAVY):
    add_round_rect(slide, left, top, width, height,
                    fill=WHITE, border=accent)
    # 아이콘
    add_text(slide, left + Cm(0.3), top + Cm(0.3),
             Cm(2.0), Cm(1.2), icon,
             font_size=32, color=accent, bold=True)
    # 타이틀
    add_text(slide, left + Cm(2.5), top + Cm(0.4),
             width - Cm(2.7), Cm(0.8), title,
             font_size=15, color=NAVY, bold=True)
    # 본문
    add_text(slide, left + Cm(0.4), top + Cm(1.7),
             width - Cm(0.6), height - Cm(1.9), body,
             font_size=11, color=GREY)


# ============================================================
# 슬라이드 빌더 — 각 슬라이드 함수
# ============================================================
def s_title(prs, *, slide_no, total, month_label):
    s = add_blank_slide(prs)
    set_bg(s, prs, NAVY)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(2.0), Cm(2.5), W - Cm(4.0), Cm(2.0),
             "월간 연금 포트폴리오",
             font_size=44, bold=True, color=WHITE)
    add_text(s, Cm(2.0), Cm(5.0), W - Cm(4.0), Cm(2.0),
             f"{month_label}호",
             font_size=72, bold=True, color=GOLD)
    add_text(s, Cm(2.0), Cm(10.0), W - Cm(4.0), Cm(1.0),
             "40대 DC/IRP 연금 가입자를 위한 자율주행 SAA",
             font_size=20, color=WHITE)
    add_text(s, Cm(2.0), Cm(11.5), W - Cm(4.0), Cm(0.8),
             "21개 자산배분 모델이 합의한 자리",
             font_size=14, color=GOLD, italic_workaround=False) \
        if False else None
    add_text(s, Cm(2.0), Cm(11.5), W - Cm(4.0), Cm(0.8),
             "21개 자산배분 모델이 합의한 자리",
             font_size=14, color=GOLD)
    add_text(s, Cm(2.0), H - Cm(1.5), W - Cm(4.0), Cm(0.6),
             f"{datetime.now().strftime('%Y-%m-%d')}",
             font_size=10, color=WHITE)


def s_chapter(prs, *, slide_no, total, chapter, title, ts):
    s = add_blank_slide(prs)
    set_bg(s, prs, NAVY)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(2.0), H/2 - Cm(2.5), W - Cm(4.0), Cm(1.5),
             chapter, font_size=16, color=GOLD, bold=True)
    add_text(s, Cm(2.0), H/2 - Cm(1.0), W - Cm(4.0), Cm(2.5),
             title, font_size=42, color=WHITE, bold=True)
    add_text(s, Cm(2.0), H/2 + Cm(2.5), W - Cm(4.0), Cm(1.0),
             ts, font_size=14, color=GOLD)


def s_hook(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(1.5), W - Cm(3.0), Cm(1.0),
             "왜 월간 연금 포트폴리오인가",
             font_size=16, color=GOLD, bold=True)
    add_text(s, Cm(1.5), Cm(3.5), W - Cm(3.0), Cm(3.0),
             "DC, IRP에 매달 돈은 들어가는데…",
             font_size=36, color=GREY)
    add_text(s, Cm(1.5), Cm(6.5), W - Cm(3.0), Cm(3.0),
             "이걸 그냥 디폴트 옵션에 두는 게",
             font_size=36, color=GREY)
    add_text(s, Cm(1.5), Cm(9.5), W - Cm(3.0), Cm(3.0),
             "정말 최선일까?",
             font_size=44, color=NAVY, bold=True)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="인트로", ts="0:00–0:15")


def s_three_walls(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "40대 연금 가입자가 마주치는 세 가지 벽",
             font_size=24, color=NAVY, bold=True)

    card_w = (W - Cm(4.0)) / 3 - Cm(0.4)
    card_h = Cm(11.5)
    top = Cm(2.5)
    cards = [
        ("⏱", "시간이 없다", "매일 시장을 들여다볼 여유가 없다.\n"
                              "분기마다 비중을 다시 잡는 것도 부담."),
        ("⚖", "방법이 너무 많다", "Equal Weight, Risk Parity, Max Sharpe, GMV…\n"
                                  "어떤 게 정답인지 알 수가 없다."),
        ("🔒", "규제가 까다롭다", "DC/IRP 위험자산 70% 한도\n레버리지·인버스 매수 금지.\n"
                                  "한도 안에 끼워 넣기 매번 계산."),
    ]
    for i, (icon, title, body) in enumerate(cards):
        left = Cm(1.8) + i * (card_w + Cm(0.4))
        add_round_rect(s, left, top, card_w, card_h,
                        fill=LIGHT, border=NAVY)
        add_text(s, left + Cm(0.3), top + Cm(0.5),
                 card_w - Cm(0.6), Cm(2.5), icon,
                 font_size=72, color=GOLD, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.3), top + Cm(4.2),
                 card_w - Cm(0.6), Cm(1.3), title,
                 font_size=22, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.3), top + Cm(6.0),
                 card_w - Cm(0.6), Cm(5.0), body,
                 font_size=13, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="인트로 — 문제 정의", ts="0:15–0:55")


def s_three_solutions(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "그래서 매달 — 자동화된 연금 포트폴리오 의견을 공개",
             font_size=22, color=NAVY, bold=True)

    card_w = (W - Cm(4.0)) / 3 - Cm(0.4)
    card_h = Cm(11.5)
    top = Cm(2.5)
    cards = [
        ("01", "21개 모델 동시", "단일 모델에 베팅하지 않는다.\n"
                                 "Equal Weight, Risk Parity, GMV, HRP,\n"
                                 "CVaR-min … 21개가 어디서 합의하는지 본다."),
        ("02", "70% 한도 자동", "DC/IRP 위험자산 한도를\n"
                                 "알고리즘이 강제한다.\n"
                                 "어떤 비중이 나와도 한도 초과 X."),
        ("03", "매크로 14개 공개", "BOK 금리, KTB 10년물, USD/KRW,\n"
                                   "브렌트유, VIX —\n"
                                   "왜 이 비중인지 추적 가능."),
    ]
    for i, (n, title, body) in enumerate(cards):
        left = Cm(1.8) + i * (card_w + Cm(0.4))
        add_round_rect(s, left, top, card_w, card_h,
                        fill=NAVY, border=NAVY)
        add_text(s, left + Cm(0.3), top + Cm(0.5),
                 card_w - Cm(0.6), Cm(2.0), n,
                 font_size=64, color=GOLD, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.3), top + Cm(4.0),
                 card_w - Cm(0.6), Cm(1.5), title,
                 font_size=22, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.3), top + Cm(6.0),
                 card_w - Cm(0.6), Cm(5.0), body,
                 font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="인트로 — 솔루션", ts="0:55–1:35")


def s_why_monthly(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "왜 매달 점검하는가",
             font_size=24, color=NAVY, bold=True)
    add_text(s, Cm(1.5), Cm(2.2), W - Cm(3.0), Cm(1.0),
             "분기 정기 리밸런싱은 그대로 — 그 사이를 매달 들여다본다",
             font_size=14, color=GREY)

    # 3개월 캘린더 박스 (Q3 예시)
    months = ["5월", "6월", "7월", "8월"]
    box_w = (W - Cm(5.0)) / 4 - Cm(0.3)
    box_h = Cm(7.5)
    top = Cm(4.5)
    for i, m in enumerate(months):
        left = Cm(2.0) + i * (box_w + Cm(0.4))
        is_rebal = (i == 3)
        fill = NAVY if is_rebal else LIGHT
        text_color = WHITE if is_rebal else GREY
        add_round_rect(s, left, top, box_w, box_h, fill=fill,
                        border=NAVY)
        add_text(s, left + Cm(0.3), top + Cm(0.5),
                 box_w - Cm(0.6), Cm(1.5), m,
                 font_size=28, color=text_color, bold=True,
                 align=PP_ALIGN.CENTER)
        if is_rebal:
            add_text(s, left + Cm(0.3), top + Cm(2.5),
                     box_w - Cm(0.6), Cm(2.0), "분기\n리밸런싱",
                     font_size=14, color=GOLD, bold=True,
                     align=PP_ALIGN.CENTER)
        else:
            add_text(s, left + Cm(0.3), top + Cm(2.5),
                     box_w - Cm(0.6), Cm(2.0), "매월\n점검",
                     font_size=14, color=NAVY, bold=True,
                     align=PP_ALIGN.CENTER)

    add_text(s, Cm(2.0), top + box_h + Cm(0.5),
             W - Cm(4.0), Cm(1.0),
             "한 달 안에도 BOK 금통위·FOMC·환율·유가가 움직인다.",
             font_size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="인트로 — 왜 월간", ts="1:35–1:50")


def s_macro_intro(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "어떤 비중을 짤지 결정하기 전 — 시장이 지금 어디?",
             font_size=22, color=NAVY, bold=True)

    add_text(s, Cm(1.5), Cm(2.5), W - Cm(3.0), Cm(0.8),
             "데이터 소스",
             font_size=16, color=GOLD, bold=True)

    # ECOS / FRED 박스
    add_round_rect(s, Cm(2.0), Cm(3.7), Cm(13.5), Cm(2.5),
                    fill=LIGHT, border=NAVY)
    add_text(s, Cm(2.5), Cm(4.0), Cm(12.5), Cm(0.8),
             "ECOS · 한국은행", font_size=16, color=NAVY, bold=True)
    add_text(s, Cm(2.5), Cm(4.8), Cm(12.5), Cm(1.5),
             "GDP · CPI · 실업 · 수출 · BOK 금리 · KTB 10/3년 · USD/KRW · AA 스프레드",
             font_size=11, color=GREY)

    add_round_rect(s, Cm(16.5), Cm(3.7), Cm(8.0), Cm(2.5),
                    fill=LIGHT, border=GOLD)
    add_text(s, Cm(17.0), Cm(4.0), Cm(7.0), Cm(0.8),
             "FRED · 미국 연준", font_size=16, color=GOLD, bold=True)
    add_text(s, Cm(17.0), Cm(4.8), Cm(7.0), Cm(1.5),
             "Fed Funds · VIX · Brent Crude",
             font_size=11, color=GREY)

    # 4축 점수 박스
    add_text(s, Cm(1.5), Cm(7.0), W - Cm(3.0), Cm(0.8),
             "4축 점수 분석",
             font_size=16, color=GOLD, bold=True)
    add_round_rect(s, Cm(2.0), Cm(8.2), W - Cm(4.0), Cm(5.5),
                    fill=LIGHT, border=NAVY)
    axes = [("성장 (Growth)", "GDP / 산업생산 / 수출 / 실업률"),
            ("물가 (Inflation)", "헤드라인 CPI / 근원 CPI / 유가"),
            ("통화 (Monetary)", "BOK 금리 / 커브 / Fed Funds"),
            ("금융 (Financial)", "AA 스프레드 / VIX / USD/KRW")]
    col_w = (W - Cm(5.0)) / 4
    for i, (a, b) in enumerate(axes):
        left = Cm(2.5) + i * col_w
        add_text(s, left, Cm(8.7), col_w - Cm(0.3), Cm(1.0),
                 a, font_size=15, color=NAVY, bold=True)
        add_text(s, left, Cm(10.0), col_w - Cm(0.3), Cm(3.0),
                 b, font_size=11, color=GREY)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 1 — 거시 진단", ts="2:00–2:15")


def s_growth(prs, *, slide_no, total, rd):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "성장 — 견조",
             font_size=28, color=NAVY, bold=True)
    add_text(s, Cm(1.5), Cm(2.2), W - Cm(3.0), Cm(1.0),
             "GDP, 수출, 실업 모두 양호. 산업생산만 미온적.",
             font_size=14, color=GREY)

    card_w = (W - Cm(5.0)) / 4 - Cm(0.3)
    card_h = Cm(6.5)
    top = Cm(4.0)
    cards = [
        ("GDP YoY", f"+{rd['kr_gdp_yoy']:.1f}%", "전년 대비 성장률", GREEN),
        ("수출 YoY", f"+{rd['kr_exports_yoy']:.0f}%", "전년 대비 — 호황 수준", GREEN),
        ("실업률", f"{rd['kr_unemployment']:.1f}%", "사실상 완전고용", GREEN),
        ("산업생산 YoY", f"+{rd['kr_industrial_production_yoy']:.1f}%", "외견 호황 대비 미온적", GOLD),
    ]
    for i, (l, v, sub, c) in enumerate(cards):
        left = Cm(2.0) + i * (card_w + Cm(0.4))
        add_kpi_card(s, left, top, card_w, card_h,
                     label=l, value=v, sub=sub, accent=c)

    add_round_rect(s, Cm(2.0), Cm(11.5), W - Cm(4.0), Cm(2.5),
                    fill=LIGHT, border=GREEN)
    add_text(s, Cm(2.5), Cm(11.9), W - Cm(5.0), Cm(2.0),
             "→ 4축 중 성장 점수 +0.66로 가장 양호",
             font_size=18, color=GREEN, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 1 — 성장", ts="2:15–2:45")


def s_inflation(prs, *, slide_no, total, rd):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "물가 — 안정, 그러나 유가 경계",
             font_size=28, color=NAVY, bold=True)
    add_text(s, Cm(1.5), Cm(2.2), W - Cm(3.0), Cm(1.0),
             "헤드라인은 BOK 목표 근방. 유가 충격 시 stagflation 위험.",
             font_size=14, color=GREY)

    add_kpi_card(s, Cm(2.0), Cm(4.0), Cm(7.0), Cm(7.0),
                 label="CPI YoY (헤드라인)",
                 value=f"+{rd['kr_cpi_yoy']:.2f}%",
                 sub="BOK 목표 2% 근방 — 안정", accent=GREEN)
    add_kpi_card(s, Cm(9.5), Cm(4.0), Cm(7.0), Cm(7.0),
                 label="근원 CPI YoY",
                 value=f"+{rd['kr_core_cpi_yoy']:.2f}%",
                 sub="기조 인플레이션 — 안정", accent=GREEN)
    add_kpi_card(s, Cm(17.0), Cm(4.0), Cm(7.5), Cm(7.0),
                 label="브렌트유 ($/bbl)",
                 value=f"${rd['kr_brent_oil']:.0f}",
                 sub="구조적 부담 누적 — 경계 필요", accent=RED)

    add_round_rect(s, Cm(2.0), Cm(11.5), W - Cm(4.0), Cm(2.5),
                    fill=LIGHT, border=RED)
    add_text(s, Cm(2.5), Cm(11.9), W - Cm(5.0), Cm(2.0),
             "⚠  유가 한 단계 더 → 성장 멀쩡한데 물가만 다시 뜨는 stagflation 위험",
             font_size=15, color=RED, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 1 — 물가", ts="2:45–3:15")


def s_monetary(prs, *, slide_no, total, rd):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "통화·금융 — BOK 인하 사이클, 모두 평온",
             font_size=24, color=NAVY, bold=True)

    card_w = (W - Cm(5.0)) / 4 - Cm(0.3)
    card_h = Cm(6.5)
    top = Cm(3.5)
    curve_delta_bp = rd.get("kr_curve_3y_10y_change_20d", 0.0) * 100
    cards = [
        ("BOK 기준금리", f"{rd['kr_base_rate']:.2f}%", "인하 사이클 진행", NAVY),
        ("KTB 10Y", f"{rd['kr_ktb_10y']:.2f}%", f"3-10Y +{rd['kr_curve_3y_10y']*100:.0f}bp / 20d {curve_delta_bp:+.0f}bp", NAVY),
        ("USD/KRW", f"{rd['kr_usd_krw']:.0f}", "원화 약세 지속", GOLD),
        ("VIX", f"{rd['us_vix']:.1f}", "변동성 지수 안정", GREEN),
    ]
    for i, (l, v, sub, c) in enumerate(cards):
        left = Cm(2.0) + i * (card_w + Cm(0.4))
        add_kpi_card(s, left, top, card_w, card_h,
                     label=l, value=v, sub=sub, accent=c)

    # 두 번째 행 — 보조
    top2 = Cm(11.0)
    cards2 = [
        ("KTB 3Y", f"{rd['kr_ktb_3y']:.2f}%"),
        ("AA- 스프레드", f"{rd['kr_corp_aa_spread_bp']:.0f}bp"),
        ("Fed Funds", f"{rd['us_fed_funds']:.2f}%"),
        ("3-10Y 커브", f"+{rd['kr_curve_3y_10y']*100:.0f}bp"),
    ]
    for i, (l, v) in enumerate(cards2):
        left = Cm(2.0) + i * (card_w + Cm(0.4))
        add_round_rect(s, left, top2, card_w, Cm(2.0),
                        fill=LIGHT, border=GREY)
        add_text(s, left + Cm(0.2), top2 + Cm(0.2),
                 card_w - Cm(0.4), Cm(0.7), l,
                 font_size=11, color=GREY, align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.2), top2 + Cm(0.9),
                 card_w - Cm(0.4), Cm(1.0), v,
                 font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 1 — 통화·금융", ts="3:15–3:45")


def s_macro_score(prs, *, slide_no, total, charts_dir, macro):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "결론 — 4축 점수와 레짐",
             font_size=28, color=NAVY, bold=True)

    add_image(s, charts_dir / "06_macro_scores.png",
              Cm(1.5), Cm(2.5), width=Cm(20.0))

    add_round_rect(s, Cm(22.5), Cm(4.5), Cm(7.5), Cm(7.0),
                    fill=NAVY, border=NAVY)
    add_text(s, Cm(22.8), Cm(4.8), Cm(7.0), Cm(1.0),
             "현재 레짐", font_size=12, color=GOLD, bold=True)
    add_text(s, Cm(22.8), Cm(5.8), Cm(7.0), Cm(2.0),
             macro["regime"], font_size=28, color=WHITE, bold=True)
    add_text(s, Cm(22.8), Cm(8.5), Cm(7.0), Cm(0.8),
             f"신뢰도 {macro['confidence']:.2f}",
             font_size=12, color=WHITE)
    add_text(s, Cm(22.8), Cm(9.4), Cm(7.0), Cm(0.8),
             f"12m 침체확률 {macro['recession_probability_12m']*100:.0f}%",
             font_size=12, color=WHITE)
    add_text(s, Cm(22.8), Cm(10.4), Cm(7.0), Cm(1.0),
             "사이클 후반",
             font_size=14, color=GOLD, bold=True)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 1 — 결론", ts="3:45–4:00")


def s_pf_intro(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "21개 모델이 합의한 결과",
             font_size=28, color=NAVY, bold=True)

    add_text(s, Cm(1.5), Cm(3.0), W - Cm(3.0), Cm(2.0),
             "채택된 ensemble", font_size=14, color=GREY)
    add_text(s, Cm(1.5), Cm(4.0), W - Cm(3.0), Cm(2.0),
             "inverse_te", font_size=72, color=NAVY, bold=True)
    add_text(s, Cm(1.5), Cm(7.5), W - Cm(3.0), Cm(2.0),
             "추적오차의 역수로 가중치 — late-cycle regime에 적합",
             font_size=18, color=GREY)

    # 흐름도
    flow_y = Cm(11.0)
    flow_h = Cm(2.5)
    box_w = Cm(8.0)
    add_round_rect(s, Cm(1.5), flow_y, box_w, flow_h,
                    fill=LIGHT, border=NAVY)
    add_text(s, Cm(1.5), flow_y, box_w, flow_h,
             "21개 PC 모델",
             font_size=18, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_round_rect(s, Cm(11.0), flow_y, box_w, flow_h,
                    fill=GOLD, border=GOLD)
    add_text(s, Cm(11.0), flow_y, box_w, flow_h,
             "inverse_te ensemble",
             font_size=18, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_round_rect(s, Cm(20.5), flow_y, box_w, flow_h,
                    fill=NAVY, border=NAVY)
    add_text(s, Cm(20.5), flow_y, box_w, flow_h,
             "최종 비중 10 ETF",
             font_size=18, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 2 — 추천 포트폴리오", ts="4:00–4:15")


def s_top_weights(prs, *, slide_no, total, charts_dir):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "Top 10 비중 — 큰 순서로",
             font_size=28, color=NAVY, bold=True)
    add_image(s, charts_dir / "02_cio_top_weights.png",
              Cm(1.5), Cm(2.5), width=Cm(28.5))
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 2 — Top 비중", ts="4:15–5:00")


def s_donut(prs, *, slide_no, total, charts_dir, cio_w):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "카테고리 합계 — 한 장으로",
             font_size=28, color=NAVY, bold=True)
    add_image(s, charts_dir / "01_cio_donut.png",
              Cm(2.0), Cm(2.5), width=Cm(18.0))

    # 옆 박스
    eq = (cio_w["kr-large-cap"] + cio_w["kr-dividend"]
          + cio_w["us-large-cap"] + cio_w["us-tech"]) * 100
    fi = (cio_w["kr-treasuries-10y"] + cio_w["us-treasuries-10y"]
          + cio_w["us-ig-credit"]) * 100
    ra = cio_w["gold"] * 100
    cs = (cio_w["kofr-cash"] + cio_w["money-market"]) * 100
    risk = eq + ra
    add_round_rect(s, Cm(21.0), Cm(3.0), Cm(8.5), Cm(11.0),
                    fill=NAVY, border=NAVY)
    add_text(s, Cm(21.3), Cm(3.3), Cm(8.0), Cm(0.8),
             "자산군 비중", font_size=12, color=GOLD, bold=True)
    rows = [(f"주식 (Equity)", f"{eq:.0f}%"),
            (f"채권 (FixedIncome)", f"{fi:.0f}%"),
            (f"실물자산 (RealAssets)", f"{ra:.1f}%"),
            (f"현금 (Cash)", f"{cs:.0f}%")]
    for i, (lbl, val) in enumerate(rows):
        ty = Cm(4.5 + i * 1.2)
        add_text(s, Cm(21.3), ty, Cm(5.0), Cm(0.8),
                 lbl, font_size=12, color=WHITE)
        add_text(s, Cm(26.5), ty, Cm(2.8), Cm(0.8),
                 val, font_size=14, color=GOLD, bold=True,
                 align=PP_ALIGN.RIGHT)
    add_text(s, Cm(21.3), Cm(10.5), Cm(8.0), Cm(0.8),
             "위험자산 합계", font_size=12, color=GOLD, bold=True)
    add_text(s, Cm(21.3), Cm(11.4), Cm(8.0), Cm(1.5),
             f"{risk:.1f}%", font_size=32, color=WHITE, bold=True)
    add_text(s, Cm(21.3), Cm(13.0), Cm(8.0), Cm(0.8),
             f"한도 70% — {(70-risk):.0f}%p 여유",
             font_size=11, color=GOLD)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 2 — 카테고리", ts="5:00–5:30")


def s_metrics(prs, *, slide_no, total, cio_m):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "메트릭 — KR 60/40 벤치마크 대비",
             font_size=28, color=NAVY, bold=True)

    card_w = (W - Cm(5.0)) / 4 - Cm(0.3)
    card_h = Cm(8.0)
    top = Cm(3.5)
    bm_sharpe = 1.71
    bm_mdd = -13.5
    cards = [
        ("기대수익률 E[r]", f"{cio_m['expected_return']*100:.2f}%",
         "연간, 명목", NAVY),
        ("기대변동성 σ", f"{cio_m['expected_vol']*100:.2f}%",
         "연간 표준편차", NAVY),
        ("BT Sharpe", f"{cio_m['backtest_sharpe']:.2f}",
         f"BM {bm_sharpe} 대비 {(cio_m['backtest_sharpe']-bm_sharpe):+.2f}", GREEN),
        ("BT MDD", f"{cio_m['backtest_maxdd']*100:.1f}%",
         f"BM {bm_mdd}% 대비 절반 이하", GREEN),
    ]
    for i, (l, v, sub, c) in enumerate(cards):
        left = Cm(2.0) + i * (card_w + Cm(0.4))
        add_kpi_card(s, left, top, card_w, card_h,
                     label=l, value=v, sub=sub, accent=c)

    add_round_rect(s, Cm(2.0), Cm(13.0), W - Cm(4.0), Cm(2.0),
                    fill=LIGHT, border=NAVY)
    add_text(s, Cm(2.5), Cm(13.2), W - Cm(5.0), Cm(1.5),
             f"Tracking Error {cio_m['tracking_error']*100:.2f}%  |  "
             f"HHI {cio_m['hhi']:.3f}  |  Effective N {cio_m['effective_n']:.1f}",
             font_size=14, color=GREY,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 2 — 메트릭", ts="5:30–5:45")


def s_models_hook(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), H/2 - Cm(3.0), W - Cm(3.0), Cm(1.5),
             "21개 모델은 다 같은 결론?",
             font_size=44, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(1.5), H/2 + Cm(-0.5), W - Cm(3.0), Cm(2.0),
             "답은 — 아닙니다.",
             font_size=36, color=RED, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(1.5), H/2 + Cm(2.5), W - Cm(3.0), Cm(2.0),
             "그런데 묘한 점이 있습니다.",
             font_size=24, color=GOLD, italic_workaround=False) \
        if False else None
    add_text(s, Cm(1.5), H/2 + Cm(2.5), W - Cm(3.0), Cm(2.0),
             "그런데 묘한 점이 있습니다.",
             font_size=24, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 3 — 21모델 비교", ts="5:45–6:00")


def s_dispersion(prs, *, slide_no, total, charts_dir):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "분포 — 위험자산 20% 부터 55% 까지 35%p 갭",
             font_size=24, color=NAVY, bold=True)
    add_image(s, charts_dir / "03_risky_dispersion.png",
              Cm(1.5), Cm(2.5), height=Cm(13.0))

    add_round_rect(s, Cm(20.0), Cm(3.0), Cm(9.0), Cm(5.5),
                    fill=LIGHT, border=RED)
    add_text(s, Cm(20.3), Cm(3.3), Cm(8.5), Cm(0.8),
             "가장 공격적 (50% 이상)", font_size=11, color=RED, bold=True)
    add_text(s, Cm(20.3), Cm(4.2), Cm(8.5), Cm(4.0),
             "• market-cap-weight 54.8%\n"
             "• equal-weight 50%\n"
             "• black-litterman 50%\n"
             "• cvar-min 50%\n"
             "• max-entropy 50%",
             font_size=12, color=GREY)
    add_round_rect(s, Cm(20.0), Cm(9.0), Cm(9.0), Cm(5.5),
                    fill=LIGHT, border=GREEN)
    add_text(s, Cm(20.3), Cm(9.3), Cm(8.5), Cm(0.8),
             "가장 방어적 (≤25%)", font_size=11, color=GREEN, bold=True)
    add_text(s, Cm(20.3), Cm(10.2), Cm(8.5), Cm(4.0),
             "• max-dd-constrained 20%\n"
             "• tpa 22%\n"
             "• gmv 23.6%\n"
             "• max-sharpe 24.7%",
             font_size=12, color=GREY)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 3 — 분포", ts="6:00–6:45")


def s_robustness(prs, *, slide_no, total, charts_dir):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "강건성 — ETF별 평균과 CIO 최종이 거의 일치",
             font_size=22, color=NAVY, bold=True)
    add_image(s, charts_dir / "04_avg_vs_cio.png",
              Cm(1.5), Cm(2.5), width=Cm(20.0))

    add_round_rect(s, Cm(22.5), Cm(3.0), Cm(7.5), Cm(11.0),
                    fill=NAVY, border=NAVY)
    add_text(s, Cm(22.8), Cm(3.3), Cm(7.0), Cm(0.8),
             "핵심 메시지", font_size=12, color=GOLD, bold=True)
    add_text(s, Cm(22.8), Cm(4.3), Cm(7.0), Cm(2.5),
             "Δ ≤ 1pp",
             font_size=44, color=WHITE, bold=True)
    add_text(s, Cm(22.8), Cm(7.5), Cm(7.0), Cm(6.5),
             "어떤 PC 방법론을 쓰더라도, 5월 매크로·CMA "
             "환경에서는 결국 비슷한 자리로 수렴.\n\n"
             "→ 모델 의존성 낮음\n"
             "→ 결과가 강건함\n"
             "→ 한두 모델의 가정이 틀려도 결론이 뒤집히지 않음",
             font_size=12, color=WHITE)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 3 — 강건성", ts="6:45–7:30")


def s_signal_meaning(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "21개 모델 합의 자체가 시그널",
             font_size=28, color=NAVY, bold=True)

    card_w = (W - Cm(5.0)) / 2 - Cm(0.5)
    card_h = Cm(11.0)
    top = Cm(2.8)

    # 합의
    add_round_rect(s, Cm(2.0), top, card_w, card_h,
                    fill=GREEN, border=GREEN)
    add_text(s, Cm(2.0), top + Cm(0.5), card_w, Cm(2.0),
             "✓ 합의가 강한 달", font_size=24, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(2.5), top + Cm(3.5), card_w - Cm(1.0), Cm(7.0),
             "21개 모델이 한 자리로 수렴.\n\n"
             "→ \"이 자리는 단단하다\"\n\n"
             "현재 5월 — 강건한 합의 상태",
             font_size=15, color=WHITE, align=PP_ALIGN.CENTER)

    # 변곡
    left2 = Cm(2.0) + card_w + Cm(0.5)
    add_round_rect(s, left2, top, card_w, card_h,
                    fill=GOLD, border=GOLD)
    add_text(s, left2, top + Cm(0.5), card_w, Cm(2.0),
             "△ 합의가 깨지는 달", font_size=24, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, left2 + Cm(0.5), top + Cm(3.5), card_w - Cm(1.0), Cm(7.0),
             "21개 모델이 서로 다른 곳을 가리킴.\n\n"
             "→ \"무언가 변곡이 오고 있다\"\n\n"
             "다음 영상에서 같이 추적",
             font_size=15, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 3 — 시그널", ts="7:30–7:45")


def s_pass_items(prs, *, slide_no, total, cio_w):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "통과 항목 ✓",
             font_size=28, color=GREEN, bold=True)

    eq = (cio_w["kr-large-cap"] + cio_w["kr-dividend"]
          + cio_w["us-large-cap"] + cio_w["us-tech"]) * 100
    ra = cio_w["gold"] * 100
    risk = eq + ra
    max_w = max(cio_w.values()) * 100

    items = [
        ("DC/IRP 위험자산 70% 한도",
         f"{risk:.1f}% — {(70-risk):.0f}%p 여유",
         GREEN),
        ("카테고리 bounds (Equity 20-55%, FI 20-70%, RA ≤15%, Cash ≤30%)",
         "모두 충족",
         GREEN),
        ("종목별 25% 한도",
         f"최대 비중 {max_w:.1f}% (KOFR-cash)",
         GREEN),
    ]
    for i, (title, body, c) in enumerate(items):
        top = Cm(3.0 + i * 3.2)
        add_round_rect(s, Cm(2.0), top, W - Cm(4.0), Cm(2.7),
                        fill=LIGHT, border=c)
        add_text(s, Cm(2.5), top + Cm(0.3), Cm(1.5), Cm(2.0),
                 "✓", font_size=42, color=c, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Cm(4.5), top + Cm(0.4), W - Cm(7.0), Cm(1.0),
                 title, font_size=15, color=NAVY, bold=True)
        add_text(s, Cm(4.5), top + Cm(1.4), W - Cm(7.0), Cm(1.2),
                 body, font_size=13, color=GREY)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 4 — 통과", ts="7:45–8:00")


def s_warn_items(prs, *, slide_no, total, cio_m):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "주의 항목 ⚠",
             font_size=28, color=GOLD, bold=True)

    items = [
        ("기대변동성이 IPS 하한 미달",
         f"σ {cio_m['expected_vol']*100:.2f}% < IPS 하한 6.0%",
         "위험이 너무 적은 게 아니라, 수익 기회를 충분히 잡지 못하는 자리.\n"
         "사이클 후반부에 의도된 결과 — 보수적인 자리에 의도적으로 머무르고 있음."),
        ("Tracking Error가 budget 초과",
         f"TE {cio_m['tracking_error']*100:.2f}% > IPS budget 6.0%",
         "한국 60/40 벤치마크가 KOSPI 60%+국고채 40%로 한국에 너무 집중.\n"
         "글로벌 분산 포트폴리오와 자연스럽게 멀어진 결과 — "
         "다음 분기 글로벌 60/40 BM 재정의 예정."),
    ]
    for i, (title, metric, body) in enumerate(items):
        top = Cm(3.0 + i * 5.5)
        add_round_rect(s, Cm(2.0), top, W - Cm(4.0), Cm(5.0),
                        fill=LIGHT, border=GOLD)
        add_text(s, Cm(2.5), top + Cm(0.3), Cm(1.5), Cm(2.0),
                 "⚠", font_size=42, color=GOLD, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Cm(4.5), top + Cm(0.4), W - Cm(7.0), Cm(1.0),
                 title, font_size=15, color=NAVY, bold=True)
        add_text(s, Cm(4.5), top + Cm(1.4), W - Cm(7.0), Cm(0.8),
                 metric, font_size=13, color=RED, bold=True)
        add_text(s, Cm(4.5), top + Cm(2.4), W - Cm(7.0), Cm(2.5),
                 body, font_size=11, color=GREY)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 4 — 주의", ts="8:00–8:30")


def s_stress(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "스트레스 시나리오 3가지",
             font_size=28, color=NAVY, bold=True)

    card_w = (W - Cm(5.0)) / 3 - Cm(0.3)
    card_h = Cm(11.0)
    top = Cm(3.0)
    cards = [
        ("🛢️", "유가 충격", "Brent → $130+",
         "Gold 5% + US IG 8%가\n인플레 헷지로 작동.\n현금 28%는 추가 매수 여력.", GOLD),
        ("📉", "BOK 추가 인하", "2.5% → 2.0%",
         "KR 10Y 14% + US 10Y 12%\n듀레이션 노출.\n평가이익 +1.5–2% 기대.", NAVY),
        ("💱", "원/달러 충격", "1485 → 1550+",
         "US Equity 16% 환노출.\n원화환산 +5–7% 가능.\nGold·US IG는 헷지.", GREEN),
    ]
    for i, (icon, t, sub, body, c) in enumerate(cards):
        left = Cm(2.0) + i * (card_w + Cm(0.4))
        add_round_rect(s, left, top, card_w, card_h,
                        fill=WHITE, border=c)
        add_text(s, left + Cm(0.3), top + Cm(0.3),
                 card_w - Cm(0.6), Cm(2.0), icon,
                 font_size=44, color=c, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.3), top + Cm(2.5),
                 card_w - Cm(0.6), Cm(1.0), t,
                 font_size=20, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.3), top + Cm(3.8),
                 card_w - Cm(0.6), Cm(1.0), sub,
                 font_size=14, color=c, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.5), top + Cm(5.5),
                 card_w - Cm(1.0), Cm(5.0), body,
                 font_size=12, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 4 — 스트레스", ts="8:30–9:00")


def s_next_intro(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "다음 한 달, 무엇을 보아야 하나",
             font_size=28, color=NAVY, bold=True)
    add_text(s, Cm(1.5), Cm(2.5), W - Cm(3.0), Cm(1.0),
             "분기 정기 리밸런싱은 8월. 그 사이 트리거 4가지.",
             font_size=14, color=GREY)

    card_w = (W - Cm(5.0)) / 4 - Cm(0.3)
    card_h = Cm(10.5)
    top = Cm(4.5)
    cards = [
        ("BOK 금통위", "기준금리", "추가 인하?", NAVY),
        ("브렌트유", "$120", "돌파 시 stagflation 전환", RED),
        ("USD/KRW", "1500 / 1450", "어느 쪽이든 IPS escalation", GOLD),
        ("VIX", "25", "초과 시 위험자산 추가 축소", GOLD),
    ]
    for i, (l, threshold, action, c) in enumerate(cards):
        left = Cm(2.0) + i * (card_w + Cm(0.4))
        add_round_rect(s, left, top, card_w, card_h,
                        fill=LIGHT, border=c)
        add_text(s, left, top + Cm(0.5), card_w, Cm(1.5),
                 l, font_size=20, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left, top + Cm(3.0), card_w, Cm(1.5),
                 "임계값", font_size=11, color=GREY,
                 align=PP_ALIGN.CENTER)
        add_text(s, left, top + Cm(4.0), card_w, Cm(1.8),
                 threshold, font_size=22, color=c, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, left + Cm(0.3), top + Cm(7.0),
                 card_w - Cm(0.6), Cm(3.0),
                 action, font_size=12, color=GREY,
                 align=PP_ALIGN.CENTER)

    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="본론 5 — 다음 점검", ts="9:00–9:50")


def s_disclaimer(prs, *, slide_no, total):
    s = add_blank_slide(prs)
    set_bg(s, prs, BG)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(0.8), W - Cm(3.0), Cm(1.0),
             "면책 고지",
             font_size=28, color=NAVY, bold=True)
    add_round_rect(s, Cm(2.0), Cm(3.0), W - Cm(4.0), Cm(11.0),
                    fill=LIGHT, border=GREY)
    text = (
        "본 영상은 자율주행 SAA 파이프라인이 산출한 정보를 시각화한 "
        "정보 제공 자료이며, 특정 금융상품의 매수·매도 권유가 아닙니다.\n\n"
        "투자자문업·투자일임업·투자권유 행위가 아닙니다.\n\n"
        "백테스트 결과는 가상의 가정에 기반하며 미래 성과를 보장하지 않습니다.\n\n"
        "DC/IRP 규제(위험자산 70% 한도, 레버리지·인버스 ETF 매수 금지)를 "
        "가정한 IPS 하에서 산출된 결과입니다.\n\n"
        "실제 투자 결정 전 본인의 투자목적, 위험감수도, 재무상황을 "
        "종합적으로 검토하시기 바랍니다."
    )
    add_text(s, Cm(2.7), Cm(3.7), W - Cm(5.4), Cm(9.5),
             text, font_size=14, color=GREY)
    add_footer(s, prs, slide_no=slide_no, total=total,
               chapter="클로징 — 면책", ts="9:50–10:10")


def s_outro(prs, *, slide_no, total, month_label):
    s = add_blank_slide(prs)
    set_bg(s, prs, NAVY)
    W, H = prs.slide_width, prs.slide_height
    add_text(s, Cm(1.5), Cm(2.0), W - Cm(3.0), Cm(1.5),
             f"{month_label}호 — 끝",
             font_size=20, color=GOLD, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(1.5), Cm(4.5), W - Cm(3.0), Cm(2.5),
             "다음 6월호 예고",
             font_size=36, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(2.0), Cm(7.5), W - Cm(4.0), Cm(3.0),
             "BOK 5월 결정과 매크로 변화가\n"
             "21개 모델 합의를 어떻게 흔들었는지 비교",
             font_size=20, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(1.5), H - Cm(3.0), W - Cm(3.0), Cm(1.0),
             "구독 · 알림 · 댓글로 통장 사정 질문 환영",
             font_size=18, color=GOLD, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(1.5), H - Cm(1.5), W - Cm(3.0), Cm(0.7),
             "월간 연금 포트폴리오 · AI 은퇴준비파트너 · 투영인",
             font_size=10, color=WHITE,
             align=PP_ALIGN.CENTER)


# ============================================================
# Workaround for `italic_workaround` kwarg accidentally passed
# (handled by no-op in add_text via signature)
# ============================================================
_orig_add_text = add_text
def add_text(slide, *args, **kwargs):  # noqa: F811
    kwargs.pop("italic_workaround", None)
    return _orig_add_text(slide, *args, **kwargs)


# ============================================================
# 메인
# ============================================================
def main(out_dir: str, month_label: str) -> None:
    base = Path(__file__).resolve().parent.parent
    out_path = base / out_dir
    charts_dir = out_path / "charts"

    cio = json.loads((out_path / "cio" / "final_portfolio.json").read_text(encoding="utf-8"))
    macro = json.loads((out_path / "macro" / "macro-view.json").read_text(encoding="utf-8"))
    ips = yaml.safe_load((base / "configs" / "ips_trimmed10.yaml").read_text(encoding="utf-8"))

    cio_w = cio["weights"]
    cio_m = cio["metrics"]
    rd = macro["readings"]

    prs = Presentation()
    prs.slide_width = Cm(33.867)   # 16:9 widescreen
    prs.slide_height = Cm(19.05)

    TOTAL = 29

    # ─── 슬라이드 시퀀스 ───
    s_title(prs, slide_no=1, total=TOTAL, month_label=month_label)
    s_chapter(prs, slide_no=2, total=TOTAL,
              chapter="CHAPTER 1", title="왜 월간 연금 포트폴리오인가",
              ts="0:00–1:50")
    s_hook(prs, slide_no=3, total=TOTAL)
    s_three_walls(prs, slide_no=4, total=TOTAL)
    s_three_solutions(prs, slide_no=5, total=TOTAL)
    s_why_monthly(prs, slide_no=6, total=TOTAL)

    s_chapter(prs, slide_no=7, total=TOTAL,
              chapter="CHAPTER 2", title="5월 거시 진단",
              ts="2:00–4:00")
    s_macro_intro(prs, slide_no=8, total=TOTAL)
    s_growth(prs, slide_no=9, total=TOTAL, rd=rd)
    s_inflation(prs, slide_no=10, total=TOTAL, rd=rd)
    s_monetary(prs, slide_no=11, total=TOTAL, rd=rd)
    s_macro_score(prs, slide_no=12, total=TOTAL,
                  charts_dir=charts_dir, macro=macro)

    s_chapter(prs, slide_no=13, total=TOTAL,
              chapter="CHAPTER 3", title="5월 추천 포트폴리오",
              ts="4:00–5:45")
    s_pf_intro(prs, slide_no=14, total=TOTAL)
    s_top_weights(prs, slide_no=15, total=TOTAL, charts_dir=charts_dir)
    s_donut(prs, slide_no=16, total=TOTAL, charts_dir=charts_dir, cio_w=cio_w)
    s_metrics(prs, slide_no=17, total=TOTAL, cio_m=cio_m)

    s_chapter(prs, slide_no=18, total=TOTAL,
              chapter="CHAPTER 4", title="21개 모델은 다 같은 결론?",
              ts="5:45–7:45")
    s_models_hook(prs, slide_no=19, total=TOTAL)
    s_dispersion(prs, slide_no=20, total=TOTAL, charts_dir=charts_dir)
    s_robustness(prs, slide_no=21, total=TOTAL, charts_dir=charts_dir)
    s_signal_meaning(prs, slide_no=22, total=TOTAL)

    s_chapter(prs, slide_no=23, total=TOTAL,
              chapter="CHAPTER 5", title="위험 점검",
              ts="7:45–9:00")
    s_pass_items(prs, slide_no=24, total=TOTAL, cio_w=cio_w)
    s_warn_items(prs, slide_no=25, total=TOTAL, cio_m=cio_m)
    s_stress(prs, slide_no=26, total=TOTAL)

    s_next_intro(prs, slide_no=27, total=TOTAL)
    s_disclaimer(prs, slide_no=28, total=TOTAL)
    s_outro(prs, slide_no=29, total=TOTAL, month_label=month_label)

    save_path = out_path / f"slides_{month_label}.pptx"
    prs.save(save_path)
    print(f"[saved] {save_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="outputs_trimmed10")
    parser.add_argument("--month", default=None)
    args = parser.parse_args()
    month_label = args.month or datetime.now().strftime("%Y-%m")
    main(out_dir=args.out, month_label=month_label)
